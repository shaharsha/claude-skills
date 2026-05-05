#!/usr/bin/env python3
"""
build-pptx.py — assemble slide PNGs + deck plan into a 16:9 PPTX via python-pptx.

Each slide gets one full-bleed image and (optionally) the speaker notes from the plan.
Slide dimensions are 13.333" × 7.5" (the standard 16:9 widescreen).

Usage:
  build-pptx.py --plan deck-plan.json --slides-dir ./slides/ --output ./output/deck.pptx
"""

from __future__ import annotations

import argparse
import glob
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print(
        "ERROR: python-pptx not installed. Run: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(2)


SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5


def find_slide_png(slides_dir: Path, slide_id: str) -> Path | None:
    matches = sorted(slides_dir.glob(f"slide-{slide_id}-*.png"))
    if not matches:
        matches = sorted(slides_dir.glob(f"slide-{slide_id}*.png"))
    return matches[0] if matches else None


def main() -> int:
    ap = argparse.ArgumentParser(description="Build a PPTX deck from slide PNGs and a plan.")
    ap.add_argument("--plan", required=True)
    ap.add_argument("--slides-dir", required=True)
    ap.add_argument("--output", "-o", required=True)
    ap.add_argument("--no-notes", action="store_true", help="Skip speaker notes")
    args = ap.parse_args()

    plan_path = Path(args.plan).resolve()
    slides_dir = Path(args.slides_dir).resolve()
    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    if not plan_path.exists():
        print(f"ERROR: plan not found: {plan_path}", file=sys.stderr)
        return 2
    if not slides_dir.is_dir():
        print(f"ERROR: slides dir not found: {slides_dir}", file=sys.stderr)
        return 2

    with plan_path.open() as f:
        plan = json.load(f)

    slides = plan.get("slides", [])
    if not slides:
        print("ERROR: deck plan has no slides", file=sys.stderr)
        return 2

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    # Use the blank layout (index 6 in default theme); some themes vary, fall back to 5
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]

    missing_count = 0
    for s in slides:
        sid = str(s["id"])
        png = find_slide_png(slides_dir, sid)
        slide = prs.slides.add_slide(blank_layout)

        # Strip default placeholders so the slide is truly blank
        for shape in list(slide.placeholders):
            sp = shape._element
            sp.getparent().remove(sp)

        if png and png.exists():
            slide.shapes.add_picture(
                str(png),
                left=Inches(0),
                top=Inches(0),
                width=Inches(SLIDE_WIDTH_IN),
                height=Inches(SLIDE_HEIGHT_IN),
            )
        else:
            missing_count += 1
            print(f"WARN: slide {sid} image missing, leaving blank", file=sys.stderr)

        if not args.no_notes:
            notes = s.get("speaker_notes", "").strip()
            if notes:
                tf = slide.notes_slide.notes_text_frame
                tf.text = notes

    prs.save(str(out_path))

    print(
        f"✓ Wrote {out_path} "
        f"({len(slides)} slides, {missing_count} missing image{'s' if missing_count != 1 else ''})",
        file=sys.stderr,
    )
    print(out_path)
    return 0


if __name__ == "__main__":
    sys.exit(main())
