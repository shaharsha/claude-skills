#!/usr/bin/env python3
"""Embed per-slide audio files into an EXISTING pptx (small speaker icon, bottom-right).

Usage: python3 add_audio.py deck.pptx audio_dir/ out.pptx
audio_dir contains slide01.mp3, slide02.mp3, ... (number = slide position, 1-based).
Missing slide numbers are skipped. Requires: pip install python-pptx
NOTE: this only EMBEDS the audio (click-to-play). Run set_autoplay.sh afterwards —
never hand-write <p:timing> XML (it corrupts the file → PowerPoint repair dialog).
"""
import re, sys, glob, os
from pptx import Presentation
from pptx.util import Inches

def main():
    src, adir, dst = sys.argv[1], sys.argv[2], sys.argv[3]
    prs = Presentation(src)
    files = {int(re.search(r"slide(\d+)\.(mp3|m4a|wav)$", f).group(1)): f
             for f in glob.glob(os.path.join(adir, "slide*.*"))
             if re.search(r"slide(\d+)\.(mp3|m4a|wav)$", f)}
    added = 0
    for i, slide in enumerate(prs.slides, start=1):
        if i not in files: continue
        slide.shapes.add_movie(files[i],
            left=prs.slide_width - Inches(0.62), top=prs.slide_height - Inches(0.62),
            width=Inches(0.5), height=Inches(0.5),  # >=0.5in: hoverable target for the media seek bar
            mime_type="audio/mpeg")
        added += 1
    prs.save(dst)
    print(f"embedded audio on {added} slides -> {dst}")

if __name__ == "__main__":
    main()
